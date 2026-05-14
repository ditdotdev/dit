"""Generate JPMC one-pager PPT in datadatdat.com brand style."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette from datadatdat.github.io/_sass/variables.scss
NAVY = RGBColor(0x0B, 0x0B, 0x53)       # $color-primary
PURPLE = RGBColor(0x36, 0x32, 0xA7)     # $color-light-primary
GRAY = RGBColor(0xA9, 0x9D, 0xBA)       # $color-gray
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4) # $color-light-gray
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x4A)  # $color-dark-gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
BORDER = RGBColor(0xCD, 0xCD, 0xCD)

LOGO_WHITE = r"c:\dev\datadatdat\datadatdat-remote-server\images\datadatdat-white.png"
LOGO_ICON_WHITE = r"c:\dev\datadatdat\datadatdat-remote-server\images\datadatdat-icon-white.png"
LOGO_ICON_BLACK = r"c:\dev\datadatdat\datadatdat-remote-server\images\datadatdat-icon-black.png"

FONT = "Source Sans Pro"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_rect(left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

def add_text(left, top, width, height, text, *, size=12, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, italic=False,
             letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

# ---------- Hero band ----------
HERO_H = Inches(1.55)
add_rect(0, 0, SW, HERO_H, NAVY)

# Decorative purple accent band
add_rect(0, HERO_H, SW, Inches(0.06), PURPLE)

# Logo wordmark on hero
logo_h = Inches(1.0)
logo_left = Inches(0.5)
logo_top = Inches(0.27)
slide.shapes.add_picture(LOGO_WHITE, logo_left, logo_top, height=logo_h)

# Tagline next to logo (right side of hero)
add_text(Inches(4.3), Inches(0.35), Inches(8.7), Inches(0.45),
         "GIT FOR DATABASES",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
add_text(Inches(4.3), Inches(0.78), Inches(8.7), Inches(0.4),
         "Your code. Your environment. Your data.",
         size=14, color=GRAY, align=PP_ALIGN.RIGHT, italic=True)
add_text(Inches(4.3), Inches(1.12), Inches(8.7), Inches(0.35),
         "Prepared for JPMorgan Chase  |  Developer Experience & Productivity",
         size=11, color=WHITE, align=PP_ALIGN.RIGHT)

# ---------- Section title ----------
TITLE_TOP = Inches(1.78)
add_text(Inches(0.5), TITLE_TOP, Inches(12.3), Inches(0.4),
         "WHY DATADATDAT FOR JPMC", size=16, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)

# ---------- Three columns ----------
COL_TOP = Inches(2.2)
COL_H = Inches(2.75)
COL_W = Inches(4.18)
COL_GAP = Inches(0.13)
COL_LEFT_0 = Inches(0.5)

columns = [
    {
        "title": "JPMC CHALLENGES",
        "accent": GRAY,
        "items": [
            "20+ developers share a single dev database — outages and restores hit the entire team.",
            "Database rehydrations and restore-from-backup are slow and operationally disruptive.",
            "Load-bearing shared DBs make developers risk-averse; meaningful changes get avoided.",
            "Current workspaces can't host containers locally — but that's actively changing.",
            "Data workflow is misaligned with the app code workflow, eroding quality and velocity.",
        ],
    },
    {
        "title": "DATADATDAT SOLUTION",
        "accent": PURPLE,
        "items": [
            "Git semantics for databases: clone, push, pull, commit, checkout, branch.",
            "Every developer gets an isolated, fully-functional database from a known commit.",
            "Delta-only versioning of the data file and container metadata — commits are fast and lightweight.",
            "Same d3 workflow whether containers run outside or inside the workspace — future-proof.",
            "Database state versioned alongside code — aligned cognitive model with app development.",
        ],
    },
    {
        "title": "OUTCOMES FOR JPMC",
        "accent": NAVY,
        "items": [
            "Outages localized to one developer — never to the whole team.",
            "Restores collapse from hours to seconds via a single d3 checkout.",
            "Developers move faster, with confidence to make meaningful DB changes.",
            "Investment carries forward as JPMC's workspace strategy evolves.",
            "Higher code quality, faster delivery, measurably stronger dev engagement.",
        ],
    },
]

for i, col in enumerate(columns):
    left = COL_LEFT_0 + (COL_W + COL_GAP) * i
    # Card background
    add_rect(left, COL_TOP, COL_W, COL_H, WHITE, line=BORDER)
    # Accent top bar
    add_rect(left, COL_TOP, COL_W, Inches(0.32), col["accent"])
    # Title in accent bar
    add_text(left + Inches(0.2), COL_TOP, COL_W - Inches(0.4), Inches(0.32),
             col["title"], size=12, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Items
    items_tb = slide.shapes.add_textbox(left + Inches(0.2), COL_TOP + Inches(0.45),
                                        COL_W - Inches(0.4), COL_H - Inches(0.6))
    tf = items_tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for j, item in enumerate(col["items"]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        # bullet glyph
        bullet = p.add_run()
        bullet.text = "■  "
        bullet.font.name = FONT
        bullet.font.size = Pt(11)
        bullet.font.color.rgb = col["accent"]
        bullet.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.name = FONT
        body.font.size = Pt(10.5)
        body.font.color.rgb = DARK_GRAY

# ---------- Database support section (peer of the three columns) ----------
DB_TOP = Inches(5.05)
DB_H = Inches(0.95)
DB_LEFT = Inches(0.5)
DB_W = Inches(12.33)

# Card background
add_rect(DB_LEFT, DB_TOP, DB_W, DB_H, WHITE, line=BORDER)
# Accent top bar (purple to differentiate from the column accents)
add_rect(DB_LEFT, DB_TOP, DB_W, Inches(0.32), PURPLE)
# Title centered in accent bar
add_text(DB_LEFT, DB_TOP, DB_W, Inches(0.32),
         "30+ DATABASES SUPPORTED",
         size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Centered database list (two lines)
db_line = ("PostgreSQL · MySQL · MariaDB · SQL Server · Oracle · MongoDB · DynamoDB · CockroachDB · "
           "TiDB · Cassandra · ClickHouse · Couchbase · CouchDB · Redis · Valkey · Dragonfly\n"
           "Elasticsearch · OpenSearch · Neo4j · ScyllaDB · SurrealDB · TigerGraph · Qdrant · "
           "Chroma · Weaviate · Typesense · Meilisearch · InfluxDB · QuestDB · NATS · StarRocks")
add_text(DB_LEFT + Inches(0.3), DB_TOP + Inches(0.34), DB_W - Inches(0.6), DB_H - Inches(0.36),
         db_line, size=10.5, color=DARK_GRAY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- Implementation strip ----------
IMPL_TOP = Inches(6.18)
IMPL_H = Inches(0.78)
add_rect(Inches(0.5), IMPL_TOP, Inches(12.33), IMPL_H, LIGHT_GRAY)

add_text(Inches(0.7), IMPL_TOP + Inches(0.05), Inches(2.4), Inches(0.3),
         "IMPLEMENTATION", size=11, bold=True, color=NAVY)

# Three pill components with arrows between
pill_y = IMPL_TOP + Inches(0.36)
pill_h = Inches(0.36)
pill_w = Inches(3.4)
pill_gap = Inches(0.35)
pill_left_0 = Inches(0.7)

components = [
    ("d3 executable on workstation", NAVY),
    ("Docker  /  Kubernetes runtime", PURPLE),
    ("Datadatdat Remote Server (commit store)", NAVY),
]
for i, (label, color) in enumerate(components):
    left = pill_left_0 + (pill_w + pill_gap) * i
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, pill_y, pill_w, pill_h)
    pill.fill.solid(); pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    pill.adjustments[0] = 0.5
    tf = pill.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # Arrow between pills
    if i < len(components) - 1:
        arrow_left = left + pill_w + Inches(0.02)
        arrow_w = pill_gap - Inches(0.04)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left,
                                       pill_y + Inches(0.08), arrow_w, Inches(0.2))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

# ---------- Brand strip (thin closer) ----------
FOOTER_TOP = Inches(7.18)
add_rect(0, FOOTER_TOP, SW, Inches(0.32), NAVY)
add_text(Inches(0.5), FOOTER_TOP, SW - Inches(1.0), Inches(0.32),
         "datadatdat  ·  Your code. Your environment. Your data.",
         size=10, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- Save ----------
out = r"c:\dev\datadatdat\datadatdat\prospects\jpmc\jpmc-onepager.pptx"
prs.save(out)
print(f"Wrote {out}")
